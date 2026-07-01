from __future__ import annotations

import openpyxl
from pptx import Presentation

from equitymind.reporting.excel import ExcelGenerator
from equitymind.reporting.powerpoint import PowerPointGenerator


def test_excel_workbook_structure(analysis_report, tmp_path):
    path = ExcelGenerator(tmp_path).write_workbook(analysis_report, filename="wb.xlsx")
    assert path.exists()
    wb = openpyxl.load_workbook(path)
    assert wb.sheetnames == ["Overview", "Ranking", "Metrics", "Portfolio", "By Asset Class"]


def test_excel_has_live_formulas(analysis_report, tmp_path):
    path = ExcelGenerator(tmp_path).write_workbook(analysis_report, filename="wb.xlsx")
    wb = openpyxl.load_workbook(path)  # keeps formulas (data_only=False)
    ranking = wb["Ranking"]
    # Reward/Risk column (I) holds a formula, not a frozen value.
    assert str(ranking["I4"].value).startswith("=IFERROR(E4/F4")
    pivot = wb["By Asset Class"]
    # Pivot-by-formula uses COUNTIF/AVERAGEIF against the Ranking sheet.
    assert "COUNTIF(Ranking!" in str(pivot["B4"].value)
    assert "AVERAGEIF(Ranking!" in str(pivot["C4"].value)


def test_excel_native_table_present(analysis_report, tmp_path):
    path = ExcelGenerator(tmp_path).write_workbook(analysis_report, filename="wb.xlsx")
    wb = openpyxl.load_workbook(path)
    assert "RankingTable" in wb["Ranking"].tables
    assert "MetricsTable" in wb["Metrics"].tables


def test_excel_without_portfolio(analysis_report, tmp_path):
    analysis_report.portfolio = None
    path = ExcelGenerator(tmp_path).write_workbook(analysis_report, filename="wb.xlsx")
    wb = openpyxl.load_workbook(path)
    assert "Portfolio" not in wb.sheetnames
    assert "Ranking" in wb.sheetnames  # ranking/pivot still present


def test_powerpoint_deck_structure(analysis_report, tmp_path):
    path = PowerPointGenerator(tmp_path).write_deck(analysis_report, filename="deck.pptx")
    assert path.exists()
    prs = Presentation(path)
    # title + ranking + portfolio + one per asset + disclaimer
    expected = 3 + len(analysis_report.assets) + 1
    assert len(prs.slides) == expected


def test_powerpoint_contains_ticker_titles(analysis_report, tmp_path):
    path = PowerPointGenerator(tmp_path).write_deck(analysis_report, filename="deck.pptx")
    prs = Presentation(path)
    text = " ".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "EquityMind" in text
    for ticker in analysis_report.assets:
        assert ticker in text


def test_powerpoint_without_charts_is_smaller(analysis_report, tmp_path):
    with_charts = PowerPointGenerator(tmp_path).write_deck(
        analysis_report, filename="a.pptx", include_charts=True
    )
    without = PowerPointGenerator(tmp_path).write_deck(
        analysis_report, filename="b.pptx", include_charts=False
    )
    assert without.stat().st_size < with_charts.stat().st_size
