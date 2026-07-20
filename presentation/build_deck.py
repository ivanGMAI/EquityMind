"""Build the EquityMind interview deck (PPTX) that mirrors the HTML mockup.

The mockup lives at https://claude.ai/code/artifact/93f50ebf-723b-4c9f-8058-14336b73ef26
and uses cqw units (1 cqw = 1% of slide width). A 16:9 slide is 13.333 in wide,
so 1 cqw = 0.13333 in and 1 cqw of font = 9.6 pt.

Charts on slides 5-9 are drawn with matplotlib in the mockup's style from REAL
EquityMind output: ``presentation/data.json`` is a saved ``/api/result`` payload
(SBER, GAZP, LKOH, GMKN, YDEX · 1y · MOEX). The straddle on slide 9 is priced
with the project's own Black-Scholes module. Slide 4 embeds the real dashboard
screenshot, slide 1 the real QR code, slide 11 the real GigaChat commentary.

To refresh the data: POST /api/analyze with the tickers above, save the
/api/result JSON over presentation/data.json, rerun this script.

Usage:  uv run python presentation/build_deck.py
Output: presentation/EquityMind.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from matplotlib import font_manager


def _register_segoe() -> None:
    """Draw charts in Segoe UI (the deck's font) when Windows fonts are
    reachable — e.g. under WSL via /mnt/c. Falls back to DejaVu silently."""
    found = False
    for name in ("segoeui.ttf", "segoeuib.ttf", "segoeuii.ttf", "seguisb.ttf"):
        p = Path("/mnt/c/Windows/Fonts") / name
        if p.exists():
            font_manager.fontManager.addfont(str(p))
            found = True
    if found:
        plt.rcParams["font.family"] = "Segoe UI"


_register_segoe()
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
IMAGES = ROOT.parent / "images"

# Real analysis payload (saved /api/result response).
R = json.loads((ROOT / "data.json").read_text())
R = R.get("result", R)

MONTHS = {"01": "янв", "02": "фев", "03": "мар", "04": "апр", "05": "май",
          "06": "июн", "07": "июл", "08": "авг", "09": "сен", "10": "окт",
          "11": "ноя", "12": "дек"}


def ru(x: float, nd: int = 1, sign: bool = False) -> str:
    """Russian number formatting: comma decimals, proper minus sign."""
    s = f"{x:+.{nd}f}" if sign else f"{x:.{nd}f}"
    return s.replace("-", "−").replace(".", ",")


def closes(history) -> dict[str, float]:
    return {p["date"]: p["close"] for p in history if p.get("close") is not None}

# ---------------------------------------------------------------- palette ----
PAPER = RGBColor(0xFB, 0xFB, 0xF9)
INK = RGBColor(0x13, 0x1A, 0x15)
INK2 = RGBColor(0x57, 0x62, 0x5B)
INK3 = RGBColor(0x8A, 0x94, 0x8D)
ACCENT = RGBColor(0x21, 0xA0, 0x38)
ACCENT_DEEP = RGBColor(0x15, 0x77, 0x2B)
ACCENT_SOFT = RGBColor(0xE6, 0xF4, 0xE9)
WARM = RGBColor(0xC2, 0x41, 0x0C)
WARM_SOFT = RGBColor(0xF8, 0xEA, 0xE2)
LINE = RGBColor(0xE3, 0xE8, 0xE3)
LINE_SOFT = RGBColor(0xEE, 0xF1, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"

W, H = 13.333, 7.5           # slide, inches
CQ = W / 100                 # 1 cqw in inches
PADX, PADY = 5.8 * CQ, 5.2 * CQ

# matplotlib hex twins
M_PAPER, M_INK, M_INK2, M_INK3 = "#FBFBF9", "#131A15", "#57625B", "#8A948D"
M_GRID, M_AXIS = "#EEF1EE", "#E3E8E3"
M_GREEN, M_GREEN_D, M_WARM = "#21A038", "#15772B", "#C2410C"


# ---------------------------------------------------------------- helpers ----
def slide_bg(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER


def rrect(slide, x, y, w, h, fill, line_color=None, line_w=1.0, radius=0.06):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    sh.adjustments[0] = radius
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def txt(
    slide, x, y, w, h, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP, spacing=1.0, font=FONT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def runs(slide, x, y, w, h, parts, size, align=PP_ALIGN.LEFT, spacing=1.0):
    """Textbox with multiple styled runs: parts = [(text, color, bold), ...]."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    for text, color, bold in parts:
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def eyebrow(slide, label, y=PADY):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(PADX), Inches(y + 0.075),
        Inches(2.6 * CQ), Pt(2.2),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    txt(slide, PADX + 2.6 * CQ + 0.13, y, 9.0, 0.3, label.upper(), 12.5,
        ACCENT_DEEP, bold=True)


def head(slide, text, y=PADY + 0.4, size=35):
    txt(slide, PADX, y, W - 2 * PADX - 1.0, 1.6, text, size, INK, bold=True,
        spacing=1.08)


def chip(slide, x, y, text, style="plain"):
    """Pill chip; returns right edge x. style: plain | acc | warm."""
    fill, color = LINE_SOFT, INK2
    if style == "acc":
        fill, color = ACCENT_SOFT, ACCENT_DEEP
    elif style == "warm":
        fill, color = WARM_SOFT, WARM
    w = 0.45 + 0.107 * len(text)
    h = 0.42
    sh = rrect(slide, x, y, w, h, fill, radius=0.5)
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(13.5)
    r.font.bold = True
    r.font.color.rgb = color
    return x + w + 0.16


def foot(slide, num):
    txt(slide, W - 1.3, H - 0.55, 1.0, 0.3, num, 12, INK3,
        align=PP_ALIGN.RIGHT)


def flow_col(slide, x, y, w, h, label, title, items, style="plain"):
    fill, border = WHITE, LINE
    dot = ACCENT
    if style == "hi":
        fill, border = ACCENT_SOFT, ACCENT
    elif style == "pain":
        fill, border = WARM_SOFT, WARM
        dot = WARM
    rrect(slide, x, y, w, h, fill, border, 1.0, radius=0.055)
    pad = 0.29
    txt(slide, x + pad, y + 0.26, w - 2 * pad, 0.3, label.upper(), 11.5, INK3,
        bold=True)
    txt(slide, x + pad, y + 0.62, w - 2 * pad, 0.75, title, 18, INK, bold=True,
        spacing=1.1)
    iy = y + 0.66 + (0.62 if len(title) > 24 else 0.38)
    for item in items:
        runs(slide, x + pad, iy, w - 2 * pad, 0.62,
             [("•  ", dot, True), (item, INK2, False)], 14.5, spacing=1.15)
        iy += 0.40 + (0.26 if len(item) > 30 else 0)


def flow_arrow(slide, x, y):
    txt(slide, x, y, 0.5, 0.5, "→", 26, INK3, align=PP_ALIGN.CENTER)


def flow_row(slide, cols, y=2.75, h=3.1):
    """cols: list of (label, title, items, style)."""
    gap = 0.55
    cw = (W - 2 * PADX - 2 * gap) / 3
    x = PADX
    for i, (label, title, items, style) in enumerate(cols):
        flow_col(slide, x, y, cw, h, label, title, items, style)
        if i < 2:
            flow_arrow(slide, x + cw + 0.03, y + h / 2 - 0.25)
        x += cw + gap


def tile(slide, x, y, w, h, value, label, v_color=INK, v_size=30):
    rrect(slide, x, y, w, h, WHITE, LINE, 1.0, radius=0.09)
    txt(slide, x + 0.24, y + 0.2, w - 0.48, 0.6, value, v_size, v_color,
        bold=True)
    txt(slide, x + 0.24, y + 0.24 + v_size / 60, w - 0.48, h - 0.6, label,
        13, INK2, spacing=1.15)


# ----------------------------------------------------------------- charts ----
def chart_style(ax):
    ax.set_facecolor(M_PAPER)
    for s in ax.spines.values():
        s.set_visible(False)
    ax.tick_params(colors=M_INK3, labelsize=9, length=0)


def make_spark(path):
    rng = np.random.default_rng(7)
    v = np.cumsum(rng.normal(0.5, 4.5, 101)) + 120
    fig, ax = plt.subplots(figsize=(13.33, 2.55), dpi=150)
    fig.patch.set_alpha(0)
    ax.set_axis_off()
    ax.fill_between(range(101), v, v.min() - 30, color=M_GREEN, alpha=0.07)
    ax.plot(v, color=M_GREEN, lw=2.2, alpha=0.5)
    ax.margins(x=0, y=0)
    fig.subplots_adjust(0, 0, 1, 1)
    fig.savefig(path, transparent=True)
    plt.close(fig)


def make_perf(path):
    """SBER vs IMOEX cumulative return over the analysis window (real data)."""
    sber = closes(R["assets"]["SBER"]["history"])
    imoex = closes(R["benchmark"]["history"])
    dates = sorted(set(sber) & set(imoex))
    a = [(sber[d] / sber[dates[0]] - 1) * 100 for d in dates]
    b = [(imoex[d] / imoex[dates[0]] - 1) * 100 for d in dates]
    n = len(dates)
    fig, ax = plt.subplots(figsize=(6.5, 3.4), dpi=200)
    fig.patch.set_facecolor(M_PAPER)
    chart_style(ax)
    ax.yaxis.set_major_formatter(lambda v, _: f"{v:+.0f}%".replace("-", "−"))
    ax.grid(axis="y", color=M_GRID, lw=1)
    ax.axhline(0, color=M_AXIS, lw=1.5)
    ax.plot(b, color=M_INK3, lw=2)
    ax.plot(a, color=M_GREEN, lw=2.2)
    for series, color, name in ((a, M_GREEN, "SBER"), (b, M_INK3, "IMOEX")):
        ax.scatter([n - 1], [series[-1]], color=color, zorder=5, s=28)
        ax.annotate(f"{name} {ru(series[-1], 1, sign=True)}%",
                    (n - 1, series[-1]), (8, 0), textcoords="offset points",
                    color=color, fontweight="bold", fontsize=10.5, va="center")
    ticks = [round(i * (n - 1) / 4) for i in range(5)]
    ax.set_xticks(ticks, [f"{MONTHS[dates[i][5:7]]} ’{dates[i][2:4]}"
                          for i in ticks])
    ax.set_xlim(0, n + 46)
    ax.set_title(f"Кумулятивная доходность · год до {dates[-1][8:10]}."
                 f"{dates[-1][5:7]}.{dates[-1][:4]} · MOEX", loc="left",
                 fontsize=10.5, color=M_INK2, fontweight=600, pad=10)
    fig.tight_layout()
    fig.savefig(path, facecolor=M_PAPER)
    plt.close(fig)


def make_var(path):
    """Histogram of real SBER daily returns with historical VaR/CVaR marks."""
    px = list(closes(R["assets"]["SBER"]["history"]).values())
    rets = np.diff(px) / np.array(px[:-1]) * 100
    tail = R["assets"]["SBER"]["metrics"]["tail_risk"]
    var, cvar = -tail["historical_var_pct"], -tail["historical_cvar_pct"]
    counts, edges = np.histogram(rets, bins=41)
    c = (edges[:-1] + edges[1:]) / 2
    width = (edges[1] - edges[0]) * 0.86
    top = counts.max() * 1.22
    fig, ax = plt.subplots(figsize=(11.8, 3.5), dpi=200)
    fig.patch.set_facecolor(M_PAPER)
    chart_style(ax)
    for x, y in zip(c, counts):
        in_tail = x <= var
        ax.bar(x, y, width=width, color=M_WARM if in_tail else M_GREEN,
               alpha=0.9 if in_tail else 0.55)
    ax.axhline(0, color=M_AXIS, lw=1.5)
    ax.axvline(var, color=M_INK, lw=1.6, ls=(0, (5, 4)))
    ax.axvline(cvar, color=M_WARM, lw=1.6, ls=(0, (2, 4)))
    ax.text(var + 0.1, top * 0.93, f"VaR 95% · {ru(var)}%", ha="left",
            color=M_INK, fontweight="bold", fontsize=11)
    ax.text(cvar - 0.1, top * 0.80, f"CVaR · {ru(cvar)}%", ha="right",
            color=M_WARM, fontweight="bold", fontsize=10.5)
    # Caption sits at the left edge: a crash-day outlier stretches the axis
    # left, leaving free space there, while VaR labels occupy the right half.
    ax.text(c[0], top * 0.93,
            "SBER · дневные доходности за год · хвост слева от VaR",
            ha="left", color=M_INK2, fontsize=10, fontweight=600)
    lo, hi = int(np.floor(c[0])), int(np.ceil(c[-1]))
    ax.set_xticks(range(lo, hi + 1),
                  [ru(t, 0, sign=True) + "%" for t in range(lo, hi + 1)])
    ax.set_yticks([])
    ax.set_ylim(0, top)
    fig.tight_layout()
    fig.savefig(path, facecolor=M_PAPER)
    plt.close(fig)


def _asset_points():
    """Per-asset (vol%, ret%) in the frontier's own coordinates.

    Prefers ``portfolio.asset_points`` (backend ships them from the same μ/Σ
    as the frontier). For older payloads, recomputes the same way: aligned
    daily returns over the common dates, arithmetic mean × 252.
    """
    ap = R["portfolio"].get("asset_points")
    if ap:
        return {t: (a["volatility_pct"], a["expected_return_pct"])
                for t, a in ap.items()}
    histories = {t: closes(a["history"]) for t, a in R["assets"].items()}
    common = sorted(set.intersection(*(set(h) for h in histories.values())))
    out = {}
    for t, h in histories.items():
        px = np.array([h[d] for d in common])
        r = np.diff(px) / px[:-1]
        out[t] = (float(r.std(ddof=1) * np.sqrt(252) * 100),
                  float(r.mean() * 252 * 100))
    return out


def make_frontier(path):
    """Real 5-asset efficient frontier, allocations and single assets."""
    p = R["portfolio"]
    fr = sorted(p["frontier"], key=lambda x: x["return_pct"])
    vertex = min(fr, key=lambda x: x["volatility_pct"])
    # Solid efficient branch (vertex and up), dashed inefficient branch below.
    upper = [x for x in fr if x["return_pct"] >= vertex["return_pct"]]
    lower = [x for x in fr if x["return_pct"] <= vertex["return_pct"]]
    fig, ax = plt.subplots(figsize=(11.8, 3.6), dpi=200)
    fig.patch.set_facecolor(M_PAPER)
    chart_style(ax)
    ax.grid(axis="y", color=M_GRID, lw=1)
    if len(lower) > 1:
        ax.plot([x["volatility_pct"] for x in lower],
                [x["return_pct"] for x in lower],
                color=M_INK3, lw=1.6, ls=(0, (5, 4)))
    ax.plot([x["volatility_pct"] for x in upper],
            [x["return_pct"] for x in upper], color=M_GREEN_D, lw=2.5)
    ax.annotate("Эффективная граница",
                (upper[-1]["volatility_pct"], upper[-1]["return_pct"]),
                (-8, 10), textcoords="offset points", ha="right",
                color=M_GREEN_D, fontweight="bold", fontsize=11)
    # single assets — same μ/Σ basis as the frontier, so the geometry closes up
    for t, (vx, vy) in _asset_points().items():
        ax.scatter([vx], [vy], s=40, color=M_INK3, alpha=0.75, lw=0, zorder=4)
        ax.annotate(t, (vx, vy), (6, -11), textcoords="offset points",
                    color=M_INK3, fontsize=9.5, fontweight="bold")
    alloc_style = {
        "min_variance": ("#2563EB", "Мин. волатильность", (12, -14)),
        "max_sharpe": (M_GREEN, "Макс. Шарп", (10, 4)),
        "risk_parity": ("#7C3AED", "Паритет риска", (-10, -6)),
        "equal_weight": ("#B45309", "Равные веса", (10, -8)),
    }
    for key, (col, name, off) in alloc_style.items():
        al = p["allocations"][key]
        x, y = al["volatility_pct"], al["expected_return_pct"]
        ax.scatter([x], [y], s=90, color=col, edgecolors=M_PAPER, lw=2,
                   zorder=5)
        ax.annotate(name, (x, y), off, textcoords="offset points",
                    color=M_INK, fontweight="bold", fontsize=10.5,
                    ha="left" if off[0] > 0 else "right")
    ax.yaxis.set_major_formatter(lambda v, _: f"{v:+.0f}%".replace("-", "−"))
    ax.xaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.set_xlabel("Риск (годовая волатильность)", color=M_INK2, fontsize=10)
    ax.set_ylabel("Доходность (среднегодовая)", color=M_INK2, fontsize=10)
    ax.set_title("5 бумаг MOEX · год · пунктир — неэффективная ветвь · веса "
                 "без ограничений — возможны короткие позиции", loc="left",
                 fontsize=10, color=M_INK2, fontweight=600, pad=8)
    ax.margins(x=0.06, y=0.16)
    fig.tight_layout()
    fig.savefig(path, facecolor=M_PAPER)
    plt.close(fig)


def straddle_quote():
    """Price a 3-month ATM straddle on SBER with the project's own module."""
    from equitymind.derivatives.black_scholes import bs_price

    spot = R["assets"]["SBER"]["metrics"]["last_price"]
    sigma = R["assets"]["SBER"]["metrics"]["volatility"]["annualized_pct"] / 100
    strike = round(spot / 10) * 10
    r = R["portfolio"]["risk_free_rate_pct"] / 100
    t = 0.25
    call = bs_price(S=spot, K=strike, T=t, r=r, sigma=sigma,
                    option_type="call")
    put = bs_price(S=spot, K=strike, T=t, r=r, sigma=sigma, option_type="put")
    return spot, strike, sigma, call + put


def make_payoff(path):
    spot, K, sigma, prem = straddle_quote()
    lo, hi = K * 0.8, K * 1.2
    s = np.array([lo, K, hi])
    pay = np.abs(s - K) - prem
    ymin, ymax = -prem * 1.5, (hi - K) - prem + 4
    fig, ax = plt.subplots(figsize=(11.8, 3.5), dpi=200)
    fig.patch.set_facecolor(M_PAPER)
    chart_style(ax)
    ax.grid(axis="y", color=M_GRID, lw=1)
    ax.fill_between([K - prem, K + prem], ymin, 0, color=M_WARM, alpha=0.08)
    ax.fill_between([lo, K - prem], 0, ymax, color=M_GREEN, alpha=0.07)
    ax.fill_between([K + prem, hi], 0, ymax, color=M_GREEN, alpha=0.07)
    ax.axhline(0, color=M_AXIS, lw=1.5)
    ax.axvline(K, color=M_INK3, lw=1.2, ls=(0, (3, 4)))
    # Inside the axes, beside the strike line — above the axes it collides
    # with the chart title.
    ax.text(K + 1.5, ymax - 1.0, f"Страйк {K:.0f}", ha="left", va="top",
            color=M_INK2, fontsize=10, fontweight=600)
    ax.plot(s, pay, color=M_INK, lw=2.6, solid_joinstyle="round")
    # Labels go into the empty side of each breakeven: the payoff line drops
    # below zero towards the strike, so left-BE labels right-up, right-BE
    # labels left-up — centred labels sit right on the diagonal line.
    for be, off, ha in ((K - prem, (12, 8), "left"), (K + prem, (-12, 8), "right")):
        ax.scatter([be], [0], s=80, color=M_GREEN, edgecolors=M_PAPER, lw=2,
                   zorder=5)
        ax.annotate(f"Безубыток {ru(be)}", (be, 0), off,
                    textcoords="offset points", ha=ha,
                    color=M_GREEN_D, fontweight="bold", fontsize=10.5)
    ax.scatter([K], [-prem], s=64, color=M_WARM, edgecolors=M_PAPER, lw=2,
               zorder=5)
    ax.annotate(f"Макс. убыток = премия ({ru(-prem)} ₽)", (K, -prem),
                (0, -20), textcoords="offset points", ha="center",
                color=M_WARM, fontweight="bold", fontsize=10.5)
    ax.set_title(f"Длинный стрэддл на SBER · спот {ru(spot)} ₽ · страйк "
                 f"{K:.0f} · 3 мес · σ = {ru(sigma * 100)}% (историческая)",
                 loc="left", fontsize=10.5, color=M_INK2, fontweight=600,
                 pad=8)
    ax.set_xlabel("Цена SBER к экспирации, ₽", color=M_INK2, fontsize=10)
    ax.set_ylabel("P/L, ₽ на акцию", color=M_INK2, fontsize=10)
    ax.set_ylim(ymin, ymax)
    fig.tight_layout()
    fig.savefig(path, facecolor=M_PAPER)
    plt.close(fig)


def make_gradient(path):
    grad = np.linspace(-1, 1, 400).reshape(1, -1)
    fig, ax = plt.subplots(figsize=(4, 0.14), dpi=150)
    ax.set_axis_off()
    from matplotlib.colors import LinearSegmentedColormap
    cmap = LinearSegmentedColormap.from_list(
        "corr", ["#C2410C", "#EDEFEC", "#15772B"])
    ax.imshow(grad, aspect="auto", cmap=cmap)
    fig.subplots_adjust(0, 0, 1, 1)
    fig.savefig(path, transparent=True)
    plt.close(fig)


# ----------------------------------------------------------------- slides ----
def s01_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(s)
    s.shapes.add_picture(str(ASSETS / "spark.png"), 0, Inches(H - 2.55),
                         Inches(W), Inches(2.55))
    runs(s, PADX, 1.95, 11.0, 1.4,
         [("Equity", INK, True), ("Mind", ACCENT, True)], 80)
    txt(s, PADX, 3.35, 7.8, 1.0,
        "От сырых котировок — до готовой аналитики. Автоматически.",
        24, INK2, spacing=1.25)
    pill = rrect(s, PADX, 4.75, 2.85, 0.55, ACCENT, radius=0.5)
    p = pill.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "mind.foodize.ru"
    r.font.name = FONT
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.color.rgb = WHITE
    txt(s, PADX + 3.15, 4.83, 4.5, 0.4, "Иван Герасимов · 2026", 16, INK2)
    qs = 13 * CQ
    rrect(s, W - PADX - qs, H - PADY - qs, qs, qs, WHITE, LINE, 1.0,
          radius=0.07)
    s.shapes.add_picture(str(IMAGES / "qr.png"),
                         Inches(W - PADX - qs + 0.12),
                         Inches(H - PADY - qs + 0.12),
                         Inches(qs - 0.24), Inches(qs - 0.24))


def s02_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(s)
    eyebrow(s, "Проблема")
    head(s, "Между котировками и выводом — часы ручной рутины")
    flow_row(s, [
        ("Вход", "Поток данных",
         ["Котировки MOEX", "Глобальные рынки", "Новостной фон"], "plain"),
        ("Вручную · часы", "Рутина аналитика",
         ["Сбор и анализ данных", "Расчёт метрик в Excel",
          "Графики и слайды руками"], "pain"),
        ("Выход", "Один обзор",
         ["К моменту готовности рынок уже изменился"], "plain"),
    ], y=2.55, h=3.35)
    foot(s, "02")


def s03_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(s)
    eyebrow(s, "Решение")
    head(s, "EquityMind проходит этот путь за минуты")
    flow_row(s, [
        ("Данные", "Рынки",
         ["Московская биржа", "Глобальные рынки",
          "Фундаментальные показатели"], "plain"),
        ("Ядро", "Аналитика",
         ["Доходность и волатильность", "Риск: VaR, CVaR, просадки",
          "Портфели и оптимизация", "Опционы и форварды"], "hi"),
        ("Результат", "Готовый материал",
         ["Комментарий ИИ-аналитика", "Презентация PowerPoint",
          "Excel с формулами", "Ежедневная сводка"], "plain"),
    ], y=2.3, h=3.9)
    foot(s, "03")


def s04_live(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(s)
    eyebrow(s, "Живой продукт")
    head(s, "Система развёрнута и работает — можно открыть прямо сейчас")
    fx, fy, fw, fh = PADX, 2.55, W - 2 * PADX, 3.55
    rrect(s, fx, fy, fw, fh, WHITE, LINE, 1.0, radius=0.035)
    bar_h = 0.42
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(fx + 0.02),
                             Inches(fy + 0.02), Inches(fw - 0.04),
                             Inches(bar_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LINE_SOFT
    bar.line.fill.background()
    bar.shadow.inherit = False
    for i in range(3):
        dot = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(fx + 0.22 + i * 0.21),
            Inches(fy + 0.02 + bar_h / 2 - 0.065), Inches(0.13), Inches(0.13))
        dot.fill.solid()
        dot.fill.fore_color.rgb = LINE
        dot.line.fill.background()
        dot.shadow.inherit = False
    url = rrect(s, fx + 0.95, fy + 0.09, 4.2, 0.28, WHITE, LINE, 0.75,
                radius=0.5)
    p = url.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "mind.foodize.ru"
    r.font.name = FONT
    r.font.size = Pt(11.5)
    r.font.bold = True
    r.font.color.rgb = INK2
    # screenshot, cropped from the bottom to fill the frame
    px, py = fx + 0.02, fy + 0.02 + bar_h
    pw, ph = fw - 0.04, fh - 0.04 - bar_h
    pic = s.shapes.add_picture(str(ASSETS / "dashboard_clean.png"),
                               Inches(px), Inches(py), width=Inches(pw))
    full_h = pw * 846 / 1600
    pic.crop_bottom = 1 - ph / full_h
    pic.height = Inches(ph)
    x = PADX
    y = 6.35
    x = chip(s, x, y, "Интерактивная панель", "acc")
    x = chip(s, x, y, "Анализ инструментов")
    x = chip(s, x, y, "Портфели")
    chip(s, x, y, "Опционная лаборатория")
    foot(s, "04")


def s05_perf(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(s)
    eyebrow(s, "Доходность и риск")
    head(s, "Каждая доходность показана вместе с её риском")
    s.shapes.add_picture(str(ASSETS / "perf.png"), Inches(PADX), Inches(2.5),
                         Inches(6.5), Inches(3.4))
    # Real SBER figures from the saved payload — must match the chart at left.
    m = R["assets"]["SBER"]["metrics"]
    cagr = m["performance"]["annualized_return_pct"]
    vol = m["volatility"]["annualized_pct"]
    sharpe = m["performance"]["sharpe"]
    mdd = m["risk"]["max_drawdown_pct"]
    tx, ty, tw, th, gap = 7.65, 2.5, 2.32, 1.6, 0.2
    tile(s, tx, ty, tw, th, ru(cagr, 1, sign=True) + "%",
         "CAGR — годовая доходность (SBER)",
         ACCENT_DEEP if cagr >= 0 else WARM)
    tile(s, tx + tw + gap, ty, tw, th, ru(vol, 1) + "%",
         "Годовая волатильность")
    tile(s, tx, ty + th + gap, tw, th, ru(sharpe, 2),
         "Коэффициент Шарпа")
    tile(s, tx + tw + gap, ty + th + gap, tw, th, ru(mdd, 1) + "%",
         "Максимальная просадка", WARM)
    x = PADX
    y = 6.35
    x = chip(s, x, y, "Sortino · Calmar · информационный коэффициент")
    chip(s, x, y, "Бета и альфа к бенчмарку (CAPM)")
    foot(s, "05")


def s06_var(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(s)
    eyebrow(s, "Хвостовые риски")
    head(s, "VaR оценивается двумя методами — и они сверяются между собой")
    s.shapes.add_picture(str(ASSETS / "var.png"), Inches(PADX), Inches(2.65),
                         Inches(W - 2 * PADX), Inches(3.5))
    # Real SBER tail-risk numbers — the same ones the histogram is drawn from.
    tail = R["assets"]["SBER"]["metrics"]["tail_risk"]
    x = PADX
    y = 6.35
    x = chip(s, x, y,
             f"VaR 95% (исторический): −{ru(tail['historical_var_pct'])}% в день",
             "warm")
    x = chip(s, x, y,
             f"CVaR — потери в хвосте: −{ru(tail['historical_cvar_pct'])}%",
             "warm")
    chip(s, x, y, f"Параметрический VaR: −{ru(tail['parametric_var_pct'])}%")
    foot(s, "06")


def s07_frontier(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(s)
    eyebrow(s, "Портфели")
    head(s, "Оптимизация находит лучшие сочетания риска и доходности")
    s.shapes.add_picture(str(ASSETS / "frontier.png"), Inches(PADX),
                         Inches(2.6), Inches(W - 2 * PADX), Inches(3.6))
    x = PADX
    y = 6.35
    x = chip(s, x, y, "Максимальный Шарп", "acc")
    x = chip(s, x, y, "Минимальная волатильность")
    x = chip(s, x, y, "Паритет риска")
    chip(s, x, y, "Равные веса")
    foot(s, "07")


def s08_corr(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(s)
    eyebrow(s, "Диверсификация")
    head(s, "Корреляции показывают, где диверсификация настоящая")
    # Real correlation matrix from the payload; the IMOEX column comes from
    # each asset's benchmark stats (correlation with the index).
    p = R["portfolio"]
    tickers = p["tickers"]
    corr = p["correlation"]
    bench = {
        t: (R["assets"][t]["metrics"].get("benchmark") or {}).get("correlation")
        for t in tickers
    }
    with_index = all(v is not None for v in bench.values())
    names = tickers + (["IMOEX"] if with_index else [])
    M = []
    for a in tickers:
        row = [corr[a][b] for b in tickers]
        if with_index:
            row.append(bench[a])
        M.append(row)
    if with_index:
        M.append([bench[t] for t in tickers] + [1.0])
    neu, grn, wrm = (237, 239, 236), (21, 119, 43), (194, 65, 12)

    def mix(v):
        tgt = grn if v >= 0 else wrm
        t = abs(v)
        return RGBColor(*(round(n + (g - n) * t) for n, g in zip(neu, tgt)))

    cell, lab, g = 0.58, 0.75, 0.045
    gx, gy = PADX, 2.85
    for j, n in enumerate(names):
        txt(s, gx + lab + j * (cell + g), gy - 0.32, cell, 0.3, n, 10.5, INK2,
            bold=True, align=PP_ALIGN.CENTER)
    for i, row in enumerate(M):
        y = gy + i * (cell + g)
        txt(s, gx, y + cell / 2 - 0.12, lab - 0.08, 0.3, names[i], 10.5, INK2,
            bold=True)
        for j, v in enumerate(row):
            x = gx + lab + j * (cell + g)
            c = rrect(s, x, y, cell, cell, mix(v), radius=0.12)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = f"{v:.2f}".replace(".", ",")
            r.font.name = FONT
            r.font.size = Pt(10.5)
            r.font.bold = True
            r.font.color.rgb = PAPER if v > 0.55 else INK
    sx = gx + lab + len(names) * (cell + g) + 0.55
    sw = W - PADX - sx
    # Takeaway computed from the data: the tightest and the loosest pair.
    pairs = [(corr[a][b], a, b)
             for i, a in enumerate(tickers) for b in tickers[i + 1:]]
    hi_v, hi_a, hi_b = max(pairs)
    lo_v, lo_a, lo_b = min(pairs)
    txt(s, sx, 3.5, sw, 1.6,
        f"Теснее всего связаны {hi_a} и {hi_b} ({ru(hi_v, 2)}) — вместе они "
        f"почти не диверсифицируют. Слабее всего — {lo_a} и {lo_b} "
        f"({ru(lo_v, 2)}): такая пара действительно снижает риск портфеля.",
        15.5, INK2, spacing=1.3)
    txt(s, sx, 5.15, 0.4, 0.3, "−1", 11, INK3)
    s.shapes.add_picture(str(ASSETS / "gradient.png"), Inches(sx + 0.42),
                         Inches(5.18), Inches(sw - 1.0), Inches(0.12))
    txt(s, sx + sw - 0.5, 5.15, 0.5, 0.3, "+1", 11, INK3)
    foot(s, "08")


def s09_payoff(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(s)
    eyebrow(s, "Производные инструменты")
    head(s, "Опционные стратегии — с точками безубыточности и границами риска")
    s.shapes.add_picture(str(ASSETS / "payoff.png"), Inches(PADX),
                         Inches(2.65), Inches(W - 2 * PADX), Inches(3.5))
    x = PADX
    y = 6.35
    x = chip(s, x, y, "Блэк–Шоулз и все греки", "acc")
    x = chip(s, x, y, "Подразумеваемая волатильность")
    chip(s, x, y, "Стрэддл · спреды · покрытый колл · защитный пут")
    foot(s, "09")


def s10_agent(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(s)
    eyebrow(s, "ИИ-аналитик")
    head(s, "Обычный ИИ выдумывает цифры — этот вызывает готовые инструменты")
    flow_row(s, [
        ("Шаг 1", "Запрос", ["«Разбери SBER за квартал»"], "plain"),
        ("Шаг 2 · агентный режим", "Агент сам вызывает инструменты",
         ["Риск и доходность", "Стресс-сценарии", "Новости и тональность"],
         "hi"),
        ("Шаг 3", "Комментарий",
         ["Обзор, тренд, риски", "Ключевые сигналы"], "plain"),
    ], y=3.0, h=3.3)
    foot(s, "10")


def s11_ai_live(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(s)
    eyebrow(s, "ИИ-аналитик · живой пример")
    head(s, "Так выглядит комментарий системы")
    cx, cy, cw, ch = PADX, 2.3, W - 2 * PADX, 3.75
    rrect(s, cx, cy, cw, ch, ACCENT_SOFT, ACCENT, 1.0, radius=0.05)
    pad = 0.4
    txt(s, cx + pad, cy + 0.3, cw - 2 * pad, 0.35,
        "🧠 Комментарий AI-аналитика — дословно из продукта", 16.5,
        ACCENT_DEEP, bold=True)
    rows = [
        ("Сводка", [
            ("Акции Сбербанка демонстрируют устойчивую нисходящую динамику: ",
             INK, False),
            ("минус 17,99%", INK, True),
            (" за полгода на фоне высокой волатильности.", INK, False)]),
        ("Тренд", [
            ("Медвежий: цена значительно ниже долгосрочных скользящих "
             "средних, RSI — в зоне перепроданности, рыночная позиция "
             "слабая.", INK, False)]),
        ("Риск", [
            ("Шарп −1,74 и Сортино −1,89 — низкое вознаграждение за взятый "
             "риск; максимальная просадка — ", INK, False),
            ("−23,84%", INK, True), (".", INK, False)]),
    ]
    ry = cy + 0.9
    for label, parts in rows:
        txt(s, cx + pad, ry + 0.04, 1.35, 0.3, label.upper(), 12,
            ACCENT_DEEP, bold=True)
        runs(s, cx + pad + 1.55, ry, cw - 2 * pad - 1.55, 0.9, parts, 16.5,
             spacing=1.25)
        ry += 0.92
    x = PADX
    y = 6.4
    x = chip(s, x, y, "Реальный отчёт по SBER · GigaChat-2-Max", "acc")
    x = chip(s, x, y, "~30 секунд на бумагу")
    chip(s, x, y, "Не является инвестрекомендацией", "warm")
    foot(s, "11")


def s12_totals(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(s)
    eyebrow(s, "Итоги")
    head(s, "Проект в четырёх цифрах")
    data = [
        ("15+", "метрик доходности и риска — от Шарпа до CVaR"),
        ("4", "стратегии оптимизации портфеля и эффективная граница"),
        ("6", "опционных стратегий с расчётом выплат и греков"),
        ("4", "формата результата: панель, PowerPoint, Excel, сводка"),
    ]
    gap = 0.24
    tw = (W - 2 * PADX - 3 * gap) / 4
    x = PADX
    for value, label in data:
        rrect(s, x, 2.9, tw, 2.35, WHITE, LINE, 1.0, radius=0.07)
        txt(s, x + 0.28, 3.2, tw - 0.56, 0.85, value, 44, ACCENT_DEEP,
            bold=True)
        txt(s, x + 0.28, 4.15, tw - 0.56, 1.0, label, 13, INK2, spacing=1.2)
        x += tw + gap
    foot(s, "12")


def s13_final(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(s)
    txt(s, 1.6, 2.35, W - 3.2, 1.0, "Аналитика, которая делает себя сама",
        44, INK, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 2.4, 3.45, W - 4.8, 0.8,
        "Живая версия открыта — данные, риск, портфели и опционы можно "
        "посмотреть в действии.", 16.5, INK2, align=PP_ALIGN.CENTER,
        spacing=1.3)
    pw = 2.85
    px = W / 2 - (pw + 0.3 + 4.6) / 2
    pill = rrect(s, px, 4.6, pw, 0.55, ACCENT, radius=0.5)
    p = pill.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "mind.foodize.ru"
    r.font.name = FONT
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.color.rgb = WHITE
    txt(s, px + pw + 0.3, 4.68, 4.8, 0.4,
        "Иван Герасимов · ivandivan2909@gmail.com", 15, INK2)


def main() -> None:
    ASSETS.mkdir(exist_ok=True)
    make_spark(ASSETS / "spark.png")
    make_perf(ASSETS / "perf.png")
    make_var(ASSETS / "var.png")
    make_frontier(ASSETS / "frontier.png")
    make_payoff(ASSETS / "payoff.png")
    make_gradient(ASSETS / "gradient.png")

    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333 in
    prs.slide_height = Emu(6858000)  # 7.5 in
    for build in (s01_title, s02_problem, s03_solution, s04_live, s05_perf,
                  s06_var, s07_frontier, s08_corr, s09_payoff, s10_agent,
                  s11_ai_live, s12_totals, s13_final):
        build(prs)
    out = ROOT / "EquityMind.pptx"
    prs.save(out)
    print(f"Saved {out} ({out.stat().st_size // 1024} KB, "
          f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
