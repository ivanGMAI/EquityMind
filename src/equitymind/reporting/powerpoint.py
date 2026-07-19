"""PowerPoint (.pptx) deck export.

Turns an :class:`~equitymind.pipeline.AnalysisReport` into a presentation-ready
16:9 deck in the project's brand colours (Sber green): a title slide, a
cross-asset ranking table, a portfolio-allocation slide, one slide per
instrument (price chart + key metrics + AI commentary), and a closing compliance
disclaimer.

The default python-pptx template's date/footer/slide-number placeholders are
stripped from the blank layout so slides carry no footer caption. Charts are
rendered to a temp directory, embedded, and cleaned up after saving.
"""

from __future__ import annotations

import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as PresentationObj  # the class, for typing
from pptx.util import Inches, Pt

from ..ai.prompts import DISCLAIMER
from ..logging_config import get_logger
from . import charts

if TYPE_CHECKING:  # avoid an import cycle; only needed for type hints
    from ..pipeline import AnalysisReport, AssetAnalysis

logger = get_logger(__name__)

# --- Brand palette (matches the web app: Sber green + teal) ------------------
_BRAND = RGBColor(0x21, 0xA0, 0x38)  # Sber green — primary accent
_BRAND_DARK = RGBColor(0x18, 0x7A, 0x2B)  # deeper green for the title band
_INK = RGBColor(0x23, 0x2A, 0x30)  # near-black body text
_MUTED = RGBColor(0x6B, 0x72, 0x80)  # grey captions
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_ROW_ALT = RGBColor(0xEC, 0xF6, 0xEE)  # light green zebra stripe

_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)

_TREND_RU = {"bullish": "восходящий", "bearish": "нисходящий", "neutral": "боковой"}
_ALLOC_RU = {
    "Equal weight": "Равные веса",
    "Minimum variance": "Минимальная волатильность",
    "Maximum Sharpe (tangency)": "Максимальный Шарп",
    "Risk parity": "Паритет риска",
}


def _fmt_pct(value: float | None) -> str:
    return "n/a" if value is None else f"{value:+.2f}%"


def _fmt_ratio(value: float | None) -> str:
    return "n/a" if value is None else f"{value:.2f}"


class PowerPointGenerator:
    """Render :class:`AnalysisReport` objects to branded .pptx decks."""

    def __init__(self, output_dir: str | Path = "reports") -> None:
        self.output_dir = Path(output_dir)

    # ------------------------------------------------------------------ public
    def write_deck(
        self,
        report: AnalysisReport,
        *,
        filename: str | None = None,
        include_charts: bool = True,
    ) -> Path:
        """Build and persist the deck; returns the .pptx path."""
        self.output_dir.mkdir(parents=True, exist_ok=True)
        stamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
        out_path = self.output_dir / (filename or f"market_intelligence_{stamp}.pptx")
        with tempfile.TemporaryDirectory() as tmp:
            chart_dir = Path(tmp) if include_charts else None
            prs = self.build_presentation(report, chart_dir=chart_dir)
            prs.save(str(out_path))  # images are embedded before the temp dir is cleaned
        logger.info("PowerPoint deck written to %s", out_path)
        return out_path

    def build_presentation(
        self, report: AnalysisReport, *, chart_dir: Path | None = None
    ) -> PresentationObj:
        """Assemble the deck in memory. If ``chart_dir`` is set, embed charts."""
        prs = Presentation()
        prs.slide_width = _SLIDE_W
        prs.slide_height = _SLIDE_H
        self._strip_footers(prs)
        self._title_slide(prs, report)
        self._ranking_slide(prs, report)
        if report.portfolio is not None:
            self._portfolio_slide(prs, report)
        for analysis in report.assets.values():
            self._asset_slide(prs, analysis, chart_dir)
        self._disclaimer_slide(prs)
        return prs

    @staticmethod
    def _strip_footers(prs: PresentationObj) -> None:
        """Remove the template's date/footer/slide-number placeholders so slides
        carry no bottom caption."""
        for layout in prs.slide_layouts:
            for ph in list(layout.placeholders):
                ph._element.getparent().remove(ph._element)

    # ------------------------------------------------------------------ slides
    def _title_slide(self, prs: PresentationObj, report: AnalysisReport) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Green brand band across the top.
        self._rect(slide, 0, 0, _SLIDE_W, Inches(2.4), _BRAND)
        self._rect(slide, 0, Inches(2.4), _SLIDE_W, Inches(0.08), _BRAND_DARK)
        self._add_text(
            slide,
            "EquityMind",
            Inches(0.9),
            Inches(0.7),
            Inches(11.5),
            Inches(1.0),
            size=46,
            bold=True,
            color=_WHITE,
        )
        self._add_text(
            slide,
            "Аналитика финансовых рынков",
            Inches(0.9),
            Inches(1.6),
            Inches(11.5),
            Inches(0.6),
            size=22,
            color=_WHITE,
        )
        self._add_text(
            slide,
            "Отчёт по количественному анализу — метрики, риск, портфель, деривативы",
            Inches(0.9),
            Inches(3.1),
            Inches(11.5),
            Inches(0.6),
            size=18,
            color=_INK,
        )
        self._add_text(
            slide,
            f"Сформировано {report.generated_at}   ·   "
            f"AI: {report.ai_provider} ({report.ai_model})   ·   "
            f"инструментов: {len(report.assets)}",
            Inches(0.9),
            Inches(3.8),
            Inches(11.5),
            Inches(0.6),
            size=13,
            color=_MUTED,
        )

    def _ranking_slide(self, prs: PresentationObj, report: AnalysisReport) -> None:
        comp = report.comparison
        if comp is None or not comp.entries:
            return
        slide = self._content_slide(prs, "Рейтинг инструментов")
        headers = ["#", "Тикер", "Доходность", "Волатильность", "Риск", "Дох./риск", "Тренд"]
        rows = [
            [
                str(e.rank),
                e.ticker,
                _fmt_pct(e.return_pct),
                _fmt_pct(e.annualized_volatility_pct),
                "n/a" if e.risk_score is None else f"{e.risk_score:.0f}",
                _fmt_ratio(e.reward_risk_ratio),
                _TREND_RU.get(e.trend, e.trend),
            ]
            for e in comp.entries
        ]
        self._add_table(slide, headers, rows, top=Inches(1.5))

    def _portfolio_slide(self, prs: PresentationObj, report: AnalysisReport) -> None:
        p = report.portfolio
        assert p is not None
        slide = self._content_slide(prs, "Портфельная аналитика")
        self._add_text(
            slide,
            f"{len(p.tickers)} инструментов · {p.observations} наблюдений · "
            f"средняя корреляция {p.average_correlation:+.2f} · "
            f"безрисковая ставка {p.risk_free_rate * 100:.1f}%",
            Inches(0.6),
            Inches(1.35),
            Inches(12),
            Inches(0.4),
            size=13,
            color=_MUTED,
        )
        headers = ["Портфель", "Ожид. доходность", "Волатильность", "Шарп"]
        rows = [
            [
                _ALLOC_RU.get(a.label, a.label),
                _fmt_pct(a.expected_return_pct),
                _fmt_pct(a.volatility_pct),
                _fmt_ratio(a.sharpe),
            ]
            for a in p.allocations.values()
        ]
        self._add_table(slide, headers, rows, top=Inches(2.0))

    def _asset_slide(
        self, prs: PresentationObj, analysis: AssetAnalysis, chart_dir: Path | None
    ) -> None:
        m = analysis.metrics
        p = m.to_payload()
        slide = self._content_slide(prs, f"{m.ticker} — {m.name}")

        if chart_dir is not None:
            try:
                fig = charts.price_chart(analysis.history)
                png = charts.save_figure(fig, chart_dir / f"{m.ticker}_price.png")
                slide.shapes.add_picture(str(png), Inches(0.5), Inches(1.6), width=Inches(7.2))
            except Exception as exc:  # pragma: no cover - chart is optional
                logger.warning("Chart failed for %s: %s", m.ticker, exc)

        perf, tail = p.get("performance", {}), p.get("tail_risk", {})
        risk, vol = p["risk"], p["volatility"]
        bench = p.get("benchmark")
        trend = _TREND_RU.get(p["trend"]["classification"], p["trend"]["classification"])
        bullets = [
            f"Тренд: {trend}   ·   Риск {risk.get('score')}/100 ({risk.get('band')})",
            f"Накопленная {_fmt_pct(p['cumulative_return_pct'])}   ·   "
            f"Годовая (CAGR) {_fmt_pct(perf.get('annualized_return_pct'))}",
            f"Волатильность {_fmt_pct(vol.get('annualized_pct'))}   ·   "
            f"Макс. просадка {_fmt_pct(risk.get('max_drawdown_pct'))}",
            f"Шарп {_fmt_ratio(perf.get('sharpe'))}   ·   Сортино {_fmt_ratio(perf.get('sortino'))}"
            f"   ·   Кальмар {_fmt_ratio(perf.get('calmar'))}",
            f"VaR {tail.get('confidence_pct', 95):g}% {_fmt_pct(tail.get('historical_var_pct'))}"
            f"   ·   CVaR {_fmt_pct(tail.get('historical_cvar_pct'))}",
        ]
        if bench:
            bullets.append(
                f"к {bench.get('benchmark')}: β {_fmt_ratio(bench.get('beta'))} · "
                f"α {_fmt_pct(bench.get('alpha_annual_pct'))} · корр {_fmt_ratio(bench.get('correlation'))}"
            )
        self._add_bullets(slide, bullets, Inches(8.0), Inches(1.6), Inches(4.9), Inches(2.6))

        commentary = analysis.commentary
        if commentary is not None:
            self._add_text(
                slide,
                "Комментарий AI-аналитика",
                Inches(8.0),
                Inches(4.3),
                Inches(4.9),
                Inches(0.4),
                size=13,
                bold=True,
                color=_BRAND,
            )
            self._add_text(
                slide,
                commentary.summary,
                Inches(8.0),
                Inches(4.7),
                Inches(4.9),
                Inches(2.4),
                size=11,
                color=_INK,
            )

    def _disclaimer_slide(self, prs: PresentationObj) -> None:
        slide = self._content_slide(prs, "Дисклеймер")
        self._add_text(
            slide,
            f"{DISCLAIMER}\n\nПоказатели рассчитаны по историческим рыночным данным и "
            "описывают прошлое. Прошлые результаты не гарантируют будущей доходности.",
            Inches(0.8),
            Inches(2.1),
            Inches(11.5),
            Inches(3.0),
            size=16,
            color=_INK,
        )

    # ------------------------------------------------------------------ helpers
    def _content_slide(self, prs: PresentationObj, title: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Thin green brand bar at the very top + green title beneath it.
        self._rect(slide, 0, 0, _SLIDE_W, Inches(0.16), _BRAND)
        self._add_text(
            slide,
            title,
            Inches(0.5),
            Inches(0.45),
            Inches(12.3),
            Inches(0.9),
            size=28,
            bold=True,
            color=_BRAND_DARK,
        )
        return slide

    @staticmethod
    def _rect(slide, left, top, width, height, color) -> None:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False

    @staticmethod
    def _add_text(slide, text, left, top, width, height, *, size=18, bold=False, color=None):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
        return box

    @staticmethod
    def _add_bullets(slide, bullets: list[str], left, top, width, height):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, text in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = f"•  {text}"
            para.font.size = Pt(12)
            para.font.color.rgb = _INK
            para.space_after = Pt(6)

    def _add_table(self, slide, headers: list[str], rows: list[list[str]], *, top) -> None:
        n_rows, n_cols = len(rows) + 1, len(headers)
        left, width = Inches(0.5), Inches(12.3)
        height = Inches(0.42) * n_rows
        table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = h
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.bold = True
            para.font.size = Pt(12)
            para.font.color.rgb = _WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = _BRAND
        for r, row in enumerate(rows, start=1):
            for c, value in enumerate(row):
                cell = table.cell(r, c)
                cell.text = str(value)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)
                para.font.color.rgb = _INK
                cell.fill.solid()
                cell.fill.fore_color.rgb = _WHITE if r % 2 else _ROW_ALT
